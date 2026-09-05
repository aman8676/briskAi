import json
import csv
import os
import shutil
from pathlib import Path

import pdfplumber
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
from PIL import Image

from bs4 import BeautifulSoup


class UnsupportedFileTypeError(Exception):
    pass


def _find_tesseract_executable(explicit_paths=None):
    """Return the Tesseract binary path for Windows and Linux installs."""
    candidates = []

    for env_name in ("TESSERACT_PATH", "TESSERACT_CMD"):
        value = os.environ.get(env_name)
        if value:
            candidates.append(value)

    if explicit_paths:
        candidates.extend(explicit_paths)

    system_path = shutil.which("tesseract")
    if system_path:
        candidates.append(system_path)

    common_windows_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        str(Path.home() / "AppData" / "Local" / "Programs" / "Tesseract-OCR" / "tesseract.exe"),
    ]
    candidates.extend(common_windows_paths)

    # Common Linux / Docker paths
    common_linux_paths = [
        "/usr/bin/tesseract",
        "/usr/local/bin/tesseract",
        "/usr/bin/tesseract-ocr",
    ]
    candidates.extend(common_linux_paths)

    for candidate in candidates:
        if not candidate:
            continue
        path = Path(candidate).expanduser()
        if path.exists() and path.is_file():
            return str(path)

    return None


def _ensure_tesseract_available():
    """Set pytesseract's Tesseract path if it is installed but not on PATH."""
    tesseract_executable = _find_tesseract_executable()
    if tesseract_executable:
        pytesseract.pytesseract.tesseract_cmd = tesseract_executable
        return True
    return False


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".txt": _extract_txt,
        ".md": _extract_txt,
        ".json": _extract_json,
        ".csv": _extract_csv,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".html": _extract_html,
        ".htm": _extract_html,
        ".jpg": _extract_image,
        ".jpeg": _extract_image,
        ".png": _extract_image,
        ".gif": _extract_image,
        ".bmp": _extract_image,
        ".tiff": _extract_image,
        ".webp": _extract_image,
    }

    if ext not in extractors:
        raise UnsupportedFileTypeError(f"Unsupported file type: {ext}")

    return extractors[ext](file_path)


def _extract_pdf(file_path: str) -> str:
    text_parts = []
    tesseract_ok = _ensure_tesseract_available()

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text_parts.append(f"--- Page {i} ---")
            page_text = ""

            # 1. Try native text extraction
            text = page.extract_text()
            if text and text.strip():
                page_text = text.strip()

            # 2. Extract tables
            tables = page.extract_tables()
            table_parts = []
            for table_index, table in enumerate(tables, start=1):
                table_parts.append(f"--- Table {table_index} ---")
                for row in table:
                    cells = [
                        str(cell).strip() if cell is not None else ""
                        for cell in row
                    ]
                    table_parts.append(" | ".join(cells))

            # 3. If native text is empty or very sparse (< 30 chars), try OCR on this page
            if (not page_text or len(page_text) < 30) and tesseract_ok:
                try:
                    # Convert page to image using pdfplumber's pypdfium2 renderer
                    page_img = page.to_image(resolution=200).original
                    ocr_text = pytesseract.image_to_string(page_img)
                    if ocr_text and ocr_text.strip():
                        print(f"[OCR] PDF Page {i}: Extracted {len(ocr_text.strip())} chars via Tesseract OCR")
                        page_text = (page_text + "\n" + ocr_text.strip()).strip()
                except Exception as ocr_err:
                    print(f"[OCR] PDF Page {i} OCR fallback failed: {ocr_err}")

            if page_text:
                text_parts.append(page_text)
            if table_parts:
                text_parts.extend(table_parts)

    return "\n".join(text_parts).strip()


def _extract_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    text_parts = []

    # ============================================================
    # EXTRACT NORMAL PARAGRAPHS
    # ============================================================
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text.strip())

    # ============================================================
    # EXTRACT TABLES
    # ============================================================
    for table_index, table in enumerate(doc.tables, start=1):

        text_parts.append(f"--- Table {table_index} ---")

        for row in table.rows:
            cells = [
                cell.text.strip()
                for cell in row.cells
            ]

            text_parts.append(" | ".join(cells))

    return "\n".join(text_parts).strip()

def _extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def _extract_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _extract_csv(file_path: str) -> str:
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(", ".join(row))
    return "\n".join(rows).strip()


def _extract_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path, data_only=True)
    text_parts = []
    for sheet in wb.worksheets:
        text_parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            line = ", ".join(str(cell) for cell in row if cell is not None)
            if line.strip():
                text_parts.append(line)
    return "\n".join(text_parts).strip()


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides, start=1):
        text_parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    text_parts.append(" | ".join(cells))
            elif shape.has_chart:
                chart = shape.chart
                try:
                    categories = list(chart.plots[0].categories)
                except Exception:
                    categories = []
                for series in chart.series:
                    values = list(series.values)
                    if categories and len(categories) == len(values):
                        pairs = ", ".join(f"{c}: {v}" for c, v in zip(categories, values))
                    else:
                        pairs = ", ".join(str(v) for v in values)
                    text_parts.append(f"{series.name}: {pairs}")
            elif shape.has_text_frame and shape.text.strip():
                text_parts.append(shape.text)
    return "\n".join(text_parts).strip()

def _extract_html(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip()


def _extract_image(file_path: str) -> str:
    """
    Extract text from images using OCR (Tesseract).
    Also extracts image metadata.
    """
    try:
        if not _ensure_tesseract_available():
            raise FileNotFoundError(
                "Tesseract OCR is not installed or not in your PATH. "
                "Install Tesseract and set TESSERACT_PATH or add it to PATH."
            )

        image = Image.open(file_path)
        text_parts = []

        # Add image metadata
        file_name = Path(file_path).name
        text_parts.append(f"--- Image File: {file_name} ---")
        text_parts.append(f"Image size: {image.size[0]}x{image.size[1]} pixels")
        text_parts.append(f"Image format: {image.format}")
        text_parts.append("--- OCR Extracted Text ---")

        # Extract text using Tesseract OCR
        extracted_text = pytesseract.image_to_string(image)

        if extracted_text.strip():
            text_parts.append(extracted_text.strip())
        else:
            text_parts.append("[No text detected in image]")

        return "\n".join(text_parts).strip()

    except Exception as e:
        raise UnsupportedFileTypeError(f"Failed to extract text from image {file_path}: {e}")
